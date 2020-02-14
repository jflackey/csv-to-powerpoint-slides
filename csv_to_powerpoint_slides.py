import csv
import time

from pptx import Presentation

# create powerpoint object
prs = Presentation()


# open csv file
with open('sample.csv') as f:
    reader = csv.reader(f)

# write each row to a new title slide
    for row in reader:
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = row[0].replace('ï»¿', '') #remove byte order mark if present
        subtitle.text = row[1]

# create unqiue timestamped powerpoint name
timestr = 'csv_to_powerpoint_slides_' + time.strftime("%m-%d_%H.%M") + '.pptx'

# same powerpoint
prs.save(timestr)
